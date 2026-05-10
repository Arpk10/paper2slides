from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def create_pptx(slide_data: dict) -> bytes:
    """
    Takes the structured JSON output from Claude and generates a PPTX file.
    Returns the PPTX file as a byte stream.
    """
    prs = Presentation()
    
    # We will just use the default 16:9 blank slide
    # Slide layouts: 0 is Title, 1 is Title and Content
    
    slides_info = slide_data.get("slides", [])
    
    for idx, slide_item in enumerate(slides_info):
        title_text = slide_item.get("title", f"Slide {idx+1}")
        bullets = slide_item.get("bullets", [])
        notes = slide_item.get("speaker_notes", "")
        
        # Determine layout
        if idx == 0:
            layout = prs.slide_layouts[0] # Title Slide
        else:
            layout = prs.slide_layouts[1] # Title and Content
            
        slide = prs.slides.add_slide(layout)
        
        # Set Title
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            
        # Set Bullets
        if idx > 0 and len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # clear default text
            
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                
        elif idx == 0 and len(slide.placeholders) > 1:
            # Subtitle for title slide
            subtitle = slide.placeholders[1]
            subtitle.text = " | ".join(bullets)
            
        # Add speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
            
    # Save to a bytes buffer
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    
    return pptx_stream.read()
